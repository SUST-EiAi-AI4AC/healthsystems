import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank layout is index 6 in default template
    blank_layout = prs.slide_layouts[6]
    
    # Core Design Tokens (HSL/RGB harmony)
    c_dark_bg = RGBColor(7, 11, 25)       # Slate Dark 900
    c_light_bg = RGBColor(248, 250, 252)   # Slate Light 50
    c_primary = RGBColor(59, 130, 246)     # Electric Blue
    c_secondary = RGBColor(139, 92, 246)   # Deep Purple
    c_accent = RGBColor(6, 182, 212)       # Soft Cyan
    c_text_dark = RGBColor(17, 24, 39)     # Slate Dark 950
    c_text_muted = RGBColor(107, 114, 128) # Gray 500
    c_white = RGBColor(255, 255, 255)
    c_card_border = RGBColor(229, 231, 235) # Gray 200
    c_red = RGBColor(239, 68, 68)          # Alert Red
    c_green = RGBColor(16, 185, 129)       # Success Green

    def set_bg(slide, rgb_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb_color

    def add_page_header(slide, section_tag, title_text):
        # 1. Light background
        set_bg(slide, c_light_bg)
        
        # 2. Section Tag Pill
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.35))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = c_primary
        tag_shape.line.fill.background()
        tag_shape.adjustments[0] = 0.5
        
        tf_tag = tag_shape.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_left = tf_tag.margin_right = tf_tag.margin_top = tf_tag.margin_bottom = 0
        p_tag = tf_tag.paragraphs[0]
        p_tag.alignment = PP_ALIGN.CENTER
        run_tag = p_tag.add_run()
        run_tag.text = section_tag
        run_tag.font.name = "Microsoft YaHei"
        run_tag.font.size = Pt(11)
        run_tag.font.bold = True
        run_tag.font.color.rgb = c_white
        
        # 3. Title Text
        txBox = slide.shapes.add_textbox(Inches(3.2), Inches(0.3), Inches(9.3), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run_title = p.add_run()
        run_title.text = title_text
        run_title.font.name = "Microsoft YaHei"
        run_title.font.size = Pt(22)
        run_title.font.bold = True
        run_title.font.color.rgb = c_text_dark

    def add_card(slide, left, top, width, height, title, lines, title_color=c_primary):
        # Card background (rounded rectangle)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = c_white
        card.line.color.rgb = c_card_border
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.04 # small corner radius
        
        # Text box overlay
        tx = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), height - Inches(0.36))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Card title
        p_title = tf.paragraphs[0]
        run_t = p_title.add_run()
        run_t.text = title
        run_t.font.name = "Microsoft YaHei"
        run_t.font.size = Pt(14)
        run_t.font.bold = True
        run_t.font.color.rgb = title_color
        
        # Card lines
        for line in lines:
            p_line = tf.add_paragraph()
            p_line.space_before = Pt(5)
            p_line.space_after = Pt(1)
            
            # Format bullets
            bullet_prefix = ""
            if not (line.startswith("1.") or line.startswith("2.") or line.startswith("3.") or line.startswith("4.") or line.startswith("5.") or line.startswith("-") or line.startswith("★")):
                bullet_prefix = "•  "
                
            parts = line.split("**")
            is_bold = False
            first_run = True
            for part in parts:
                run_p = p_line.add_run()
                if first_run and bullet_prefix:
                    run_p.text = bullet_prefix + part
                    first_run = False
                else:
                    run_p.text = part
                run_p.font.name = "Microsoft YaHei"
                run_p.font.size = Pt(10.5)
                run_p.font.color.rgb = c_text_dark if not is_bold else title_color
                run_p.font.bold = is_bold
                is_bold = not is_bold

    # ----------------------------------------------------
    # SLIDE 1: Cover Page (Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, c_dark_bg)
    
    # Geometric decorative block
    decor = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.0))
    decor.fill.solid()
    decor.fill.fore_color.rgb = c_primary
    decor.line.fill.background()
    
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "心晴AI — 智能心理检测管理平台评估报告"
    run1.font.name = "Microsoft YaHei"
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = c_white
    
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = "融合客观生理监测、多模态AI分析与即时疗愈的数字心理健康全面评估"
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(18)
    run2.font.color.rgb = c_accent
    
    # Meta Info
    txMeta = slide.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11.0), Inches(1.5))
    tf_m = txMeta.text_frame
    tf_m.word_wrap = True
    p_m = tf_m.paragraphs[0]
    p_m.space_before = Pt(20)
    run_m = p_m.add_run()
    run_m.text = "体验入口：z.playe.top (官网) | h.playe.top (H5) | nwpuhs.cn (医生管理后台)\n项目归属：西北工业大学智能心理健康管理团队 (SUST-EiAi-AI4AC)"
    run_m.font.name = "Microsoft YaHei"
    run_m.font.size = Pt(12)
    run_m.font.color.rgb = c_text_muted

    # ----------------------------------------------------
    # SLIDE 2: 一、需求洞察 - 痛点与解决方案
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "一、需求洞察", "痛点与解决方案 (Pain Points & Solutions)")
    
    left_lines = [
        "1. **评估手段单一且主观偏见**：传统心理量表完全依赖用户自我填写，缺乏专业规范的数字化分析，结果容易偏差。",
        "2. **生理与心理数据割裂孤立**：心率、睡眠、压力等腕表生理指标，同主观心理问卷没有关联比照，无法做到多模态诊断。",
        "3. **瞬时情绪起伏难以捕捉**：心理状态变化剧烈，传统的月度量表无法记录日常遭遇特定事件时的生态瞬时情绪（EMA）。",
        "4. **依从性弱导致记录不连续**：日常填报要求高，如果页面复杂、没有系统智能推送，患者极易漏填，导致数据断档。",
        "5. **医患双向互动存在断点**：临床上，医生在诊室外难以追踪出院患者或居家患者的真实日常状态，高危危机干预滞后。"
    ]
    right_lines = [
        "1. **多模态标准化测评矩阵**：集成 SDS、SAS、SCL-90 等 10+ 种权威量表，自动生成评分和个性化医学干预意见。",
        "2. **Garmin 生理同步交叉验证**：打通 Garmin Connect API，每小时无感拉取并存储心率、深浅睡眠、SpO2及 HRV 压力值。",
        "3. **生态瞬时评估 (EMA) 交互**：通过 16 种生动表情与 1-5 级强度拉条，一秒记录瞬时情绪，输出情绪变化波形图。",
        "4. **智能推送与极简 1分钟填报**：集成极光推送与手机厂商离线通道进行填报提醒，极简流程将记录耗时缩短至 1 分钟。",
        "5. **LayuiMini 架构医生工作台**：打通前后台数据，支持高危患者红色标注预警、AI会话情感识别与危机警报介入。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.6), "传统心理关怀痛点", left_lines, title_color=c_red)
    add_card(slide, Inches(6.8), Inches(1.3), Inches(5.6), Inches(5.6), "心晴AI 多模态解决方案", right_lines, title_color=c_green)

    # ----------------------------------------------------
    # SLIDE 3: 二、团队及产品介绍 - 项目简介（名称/定位/价值）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "二、产品介绍", "项目简介：产品定位与核心价值")
    
    c1 = [
        "**产品名称**：心晴AI - 智能心理检测辅助治疗数字平台 (Health Care System)",
        "**平台组成**：\n- 官网站: **z.playe.top**\n- 移动端: **h.playe.top**\n- 医生后台: **nwpuhs.cn**",
        "**开源属性**：GitHub 开源代码库，拥有完整的前端跨平台 App、Spring Boot 后端、Garmin 同步脚本及 SQL。"
    ]
    c2 = [
        "**多模态融合评估平台**：打破传统单一主观量表局限，结合智能手环客观体征、主观测评、房树人绘画和情绪日记。",
        "**医生临床工作台**：为咨询师设计 LayuiMini 管理端，实现患者腕表数据分析、7维量表分析、HTP绘画审核与情感提取。"
    ]
    c3 = [
        "**客观量化身心状态**：打通手环客观数据与心理主观问卷，多角度交叉验证，保证心理风险评估的科学可靠性。",
        "**降低自我调节门槛**：内置共情 LLM 伴聊与多巴胺跑酷游戏，打造‘评估-疗愈-追踪-干预’轻量化心理关怀闭环。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "一、系统组成与简介", c1)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "二、产品定位", c2)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "三、核心价值", c3)

    # ----------------------------------------------------
    # SLIDE 4: 二、团队及产品介绍 - 用户画像（年龄/身份/使用场景)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "二、产品介绍", "目标用户画像与核心使用场景")
    
    p1 = [
        "**典型人群**：抑郁/焦虑症人群、高压学生与白领群体，年龄多集中在 **18-35岁** 的亚健康人群。",
        "**使用场景**：\n- 每日早晚佩戴 Garmin 手环，实现生理数据自动后台同步。\n- 每日 20:00 后完成食欲、睡眠、精力极简日常状态填报。\n- 情绪低落时，通过语音/文字与 AI 心理助手‘心晴’聊天，或在疗愈空间玩敏捷跑酷游戏排解压力。"
    ]
    p2 = [
        "**典型人群**：精神科执业医生、专业心理咨询师、学校或机构心理科专员。",
        "**使用场景**：\n- 在 Web 后台（nwpuhs.cn）分页检索并管理名下的关联患者，生成邀请码绑定患者。\n- 查看高危高分被标成红色的异常心理量表，分析房树人（HTP）AI 评语和画作大图。\n- 审查患者 AI 聊天情感度，输入诊疗意见发回移动端。"
    ]
    p3 = [
        "**典型机构**：高校校医院（如**西北工业大学校医院**）、企事业单位工会、心理数字疗法临床试验中心。",
        "**使用场景**：\n- 招募测试受试者，线下签署知情同意书，分配唯一病例号（如 SUB-001），进行 15天 闭环临床随访。\n- 结合 **近红外脑功能成像检测 (fNIRS)** 采集基准脑活动。\n- 监督手环同步和每日填报，获取多模态真实数据库。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "患者端：18-35岁青年群体", p1)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "医生端：心理师及临床医生", p2)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "机构端：校医院与科研测试团队", p3)

    # ----------------------------------------------------
    # SLIDE 5: 二、团队及产品介绍 - 核心及特色功能
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "二、产品介绍", "系统核心及特色功能矩阵")
    
    f1 = [
        "**10+套专业量表库**：支持 SCL-90（90题精神体检）、OCEAN 大五人格、PSQI 睡眠质量等权威工具。",
        "**临床评估金标准**：提供医生推荐的 SDS、SAS、PSS 每两周定期评估，自动评分与生成趋势走势图。"
    ]
    f2 = [
        "**绘画投射表达**：基于流畅手绘板，让用户手绘房屋、树木和人，画图一键上传至文件服务器。",
        "**AI 视觉量化报告**：从构图饱满度、色彩冷暖、笔触力度等 6 维度解析，识别潜在抑郁风险并归档记录。"
    ]
    f3 = [
        "**手环传感器同步**：无感同步心率（实时/静息）、睡眠分期（深/浅/REM）、步数卡路里、血氧（SpO2）、全天压力。",
        "**6维雷达健康报告**：整合生理-量表-日常填报数据，生成多维雷达健康报告，支持医生一键意见输入与 PDF 导出。"
    ]
    f4 = [
        "**心晴 AI 伴聊助手**：对接 MiniMax-M2.5 共情安抚模型，采用打字机延迟输出，提供贴心温暖对话陪伴。",
        "**SenseVoice 语音录音**：集成 FunAudioLLM/SenseVoiceSmall 模型，长按语音实时转写，解决键盘打字抗拒感。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.4), Inches(5.6), Inches(2.6), "1. 心理量表测评中心", f1)
    add_card(slide, Inches(6.8), Inches(1.4), Inches(5.6), Inches(2.6), "2. 房树人 (HTP) AI 绘画测评", f2)
    add_card(slide, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6), "3. Garmin手环数据与雷达健康报告", f3)
    add_card(slide, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.6), "4. AI 伴聊与 SenseVoice 语音输入", f4)

    # ----------------------------------------------------
    # SLIDE 6: 二、团队及产品介绍 - 竞品调研（创新差异化）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "二、产品介绍", "竞品调研与心晴AI创新差异化")
    
    comp = [
        "1. **传统心理测评 App**：\n- **局限**：纯靠用户主观打分，容易产生防御心理导致隐瞒或漏填；缺乏客观运动及睡眠数据验证，评估手段单一割裂，无互动疗愈。",
        "2. **手环及日常生理追踪软件 (Connect/Keep)**：\n- **局限**：仅从生理角度收集心率和睡眠时长，缺乏深层心理健康分析工具，无法应对抑郁/焦虑症的筛查及干预。",
        "3. **传统心理面诊及干预**：\n- **局限**：费用高、门槛重，在离开诊所后医生对患者的居家追踪完全空白，难以做到长周期实时管理。"
    ]
    diff = [
        "1. **多模态客观身心数据融合 (首创)**：\n- **优势**：打破信息壁垒，将手环客观睡眠/HRV压力指标同心理自评量表、EMA瞬时脸谱和房树人手绘AI解析融合，交叉验证评估。",
        "2. **24h 无障碍主动伴聊疗愈**：\n- **优势**：共情陪伴 LLM 配合打字机流式渲染，结合高精度 SenseVoiceSmall 语音识别转写，为低能量患者提供极低交互门槛的自我宣泄空间。",
        "3. **直通医生端的协同干预通道**：\n- **优势**：利用医生后台（nwpuhs.cn）将患者脱敏生理心理轨迹直接呈报给主治医生，高危标注和危机词自动警报为临床诊疗提供依据。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.6), "市场竞品分析与瓶颈", comp, title_color=c_text_muted)
    add_card(slide, Inches(6.8), Inches(1.3), Inches(5.6), Inches(5.6), "心晴AI 核心差异化创新优势", diff, title_color=c_primary)

    # ----------------------------------------------------
    # SLIDE 7: 二、团队及产品介绍 - 评分亮点
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "二、产品介绍", "系统评分亮点与优势评估")
    
    left = [
        "1. **极佳的可视化表现力**：ECharts 在 H5 端和 Web 后台提供一致的可视化图表，包括 6 维健康雷达图、深浅睡眠堆叠柱状图、HRV 全天压力趋势图、情绪日记效价折线图，复杂数据清晰直观。",
        "2. **厂商推送与督导闭环**：通过 JPush 极光推送和手机硬件级厂商离线通知，配合晚上 21:00 定时填报提醒。在线下结合西工大校医院临床 23:00 微信人工督导，成功保证了 15天 追踪记录依从率达 **90% 以上**。",
        "3. **客观数据无感采编**：手环通过 Garmin API 接口静默自动读取和上传数据，患者除极简填报外无需任何繁琐操作，确保了连续、完整的临床研究脱敏生理数据库沉淀。"
    ]
    
    # Custom dashboard construction for the right side
    card_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5), Inches(5.4), Inches(5.2))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = c_white
    card_right.line.color.rgb = c_card_border
    card_right.line.width = Pt(1.5)
    card_right.adjustments[0] = 0.04
    
    # Big score circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), Inches(1.9), Inches(1.5), Inches(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = c_secondary
    circle.line.fill.background()
    
    tf_c = circle.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    run_c = p_c.add_run()
    run_c.text = "95分"
    run_c.font.name = "Microsoft YaHei"
    run_c.font.size = Pt(24)
    run_c.font.bold = True
    run_c.font.color.rgb = c_white
    
    # Text next to the circle
    tx_score = slide.shapes.add_textbox(Inches(9.1), Inches(2.0), Inches(3.0), Inches(1.2))
    tf_s = tx_score.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.space_before = Pt(8)
    run_s1 = p_s.add_run()
    run_s1.text = "心晴AI 综合评价得分\n"
    run_s1.font.name = "Microsoft YaHei"
    run_s1.font.size = Pt(15)
    run_s1.font.bold = True
    run_s1.font.color.rgb = c_text_dark
    run_s2 = p_s.add_run()
    run_s2.text = "融合生理与心理评估的杰出代表"
    run_s2.font.name = "Microsoft YaHei"
    run_s2.font.size = Pt(11)
    run_s2.font.color.rgb = c_text_muted
    
    # Bullet points on metrics below
    tx_metrics = slide.shapes.add_textbox(Inches(7.3), Inches(3.6), Inches(4.8), Inches(2.8))
    tf_m = tx_metrics.text_frame
    tf_m.word_wrap = True
    
    metrics = [
        "**专业评估度**：★★★★★ (SCL-90+多维量表+房树人分析)",
        "**技术集成度**：★★★★★ (Connect API + 极光 + 双 AI 大模型)",
        "**页面体验度**：★★★★☆ (3D宣传官网 + 玻璃拟态卡片)",
        "**临床落地值**：★★★★★ (西工大校医院15天入组临床测试)"
    ]
    for m in metrics:
        p_m = tf_m.add_paragraph() if tf_m.paragraphs[0].text else tf_m.paragraphs[0]
        p_m.space_before = Pt(8)
        parts = m.split("**")
        is_bold = False
        for part in parts:
            run_pm = p_m.add_run()
            run_pm.text = part
            run_pm.font.name = "Microsoft YaHei"
            run_pm.font.size = Pt(11.5)
            run_pm.font.color.rgb = c_text_dark if not is_bold else c_primary
            run_pm.font.bold = is_bold
            is_bold = not is_bold

    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), "评分亮点说明", left)

    # ----------------------------------------------------
    # SLIDE 8: 三、技术解析及亮点 - 技术架构（前端/后端/AI能力）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "三、技术解析", "三层分离多模态技术架构")
    
    tech_front = [
        "**跨平台框架**：uni-app 框架编译，支持打包为 iOS App、Android 原生 App 及适配各手机端浏览器的 H5 页面。",
        "**UI组件与图形**：采用 uView UI 库定制卡片及表单，整合 lime-echart 实现在移动端对 ECharts 进行无缝渲染渲染。",
        "**宣传官网 (3D)**：welcome.html 使用 Outfit 字族及 CSS 3D 卡片视差动画，搭载响应式折线及雷达图。"
    ]
    tech_back = [
        "**后端核心**：Spring Boot 2.3.4 后端微服务框架，以 MyBatis 实现 ORM 关系映射，MySQL 8.0 存储身心数据。",
        "**安全认证与连接**：集成 Apache Shiro + JWT 双重校验，对用户隐私进行 RSA/AES 双向加解密，配合 Druid 数据源监控慢查询。",
        "**消息推送总线**：极光推送 JPush SDK + 各大手机厂商原生推送通道，解决 Android 后台静默离线推送难题。"
    ]
    tech_ai = [
        "**生理API数据流**：定时 Python 3.10 脚本自动对接 Garmin Connect API 同步数据，清洗后通过 HTTP POST 提交后端 API 写入数据库。",
        "**大模型共情陪伴**：ModelScope API 接入 MiniMax-M2.5 共情陪伴大模型，使用特定 System Prompt 引导进行对话。",
        "**SenseVoice 识别**：对接 SiliconFlow API 调取 FunAudioLLM/SenseVoiceSmall，实现极低延迟音频转写。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "前端表现层 (uni-app)", tech_front)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "后端持久层 (Spring Boot)", tech_back)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "AI与第三方集成层", tech_ai)

    # ----------------------------------------------------
    # SLIDE 9: 三、技术解析及亮点 - UI设计（前端/可视化)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "三、技术解析", "UI设计美学与多维度数据可视化")
    
    ui1 = [
        "**极致炫酷 3D 视觉**：welcome.html 官方宣传页采用极简深色科技风 background-color: #070b19 配合 HSL 渐变与 Outfit 现代无衬线字族，使用 CSS 3D 实现卡片旋转轮播，首创 ECharts 雷达图无感过渡，给受众带来极致的 premium 第一感官享受。"
    ]
    ui2 = [
        "**移动触控规范优化**：依据移动端交互规范设定最小 44x44px 交互触控响应区，解决错选问题。日常状态填报设计大尺寸 EMA 表情选框。加入水平滑动锁定机制，避免 scroll-view 水平拖拽同 swiper 混淆，体验流畅。"
    ]
    ui3 = [
        "**多维度可视化大屏**：",
        "- **综合雷达图**：以雷达直观比较压力、运动、睡眠及 SDS/SAS 等 6 个维度的健康平衡状况。",
        "- **深浅睡眠堆叠柱状图**：展示总睡眠时长与深浅睡眠及 REM 的睡眠结构分布，支撑精准评估睡眠效率。",
        "- **全天压力曲线**：读取佳明手环上传的 HRV 压力数据，按小时渲染全天压力起伏，直观显现全天身体负荷。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "1. 3D炫酷宣传官网 (welcome)", ui1)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "2. 移动端 H5 交互设计", ui2)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "3. 专业可视化图表设计", ui3)

    # ----------------------------------------------------
    # SLIDE 10: 三、技术解析及亮点 - 后端服务（API/数据库)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "三、技术解析", "后端 RESTful 服务与高性能数据库设计")
    
    left_backend = [
        "1. **RESTful 接口架构**：后端以统一规范的控制器（Controllers）划分业务逻辑，集成 Swagger-ui / Knife4j 实现在线调试与文档自动生成，全面覆盖用户认证、手环健康数据、量表答题和 AI 会话分析接口。",
        "2. **异步高并发任务优化**：引入 `ExecutorConfig` 与 `AsyncConfig` 对 Spring 线程池调优。针对上传音频 (SpeakAnytime)、视频抑郁推理或近红外基准比对等高能耗 Python 脚本调用任务，开启异步执行 (`Future<String>`)，保证主线程瞬间响应，提升高吞吐能力。",
        "3. **安全过滤与会话鉴权**：在 Shiro 过滤器中搭载自定义 JWT 过滤器，验证 Request Headers 中的 Token 签名与过期机制，防止伪造身份读取敏感医学记录。"
    ]
    right_backend = [
        "1. **MySQL 8.0 规范化建表**：表结构包含用户信息表、生理心率/睡眠/步数/血氧表、量表提交表、EMA 情绪明细表和 AI 聊天日志表等。对常用字段（如时间戳、病例ID）建立复合索引以提升查询速度。",
        "2. **外键约束级联控制**：对关联会话、量表记录建立 `ON DELETE CASCADE` 外键关联，支持医生或用户批量清理误删或废弃消息，规避数据库孤立脏数据产生。",
        "3. **Druid 数据库连接池**：集成阿里云 Druid 连接池并启动 SQL 监控，限制最大活动连接数与超时校验，规避数据库死锁与连接泄漏隐患。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "后端 RESTful 服务与异步多线程调度", left_backend)
    add_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "MySQL 数据库设计与持久层运维监控", right_backend)

    # ----------------------------------------------------
    # SLIDE 11: 三、技术解析及亮点 - AI能力（模型/RAG检索/Agent工作流程）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "三、技术解析", "深度 AI 能力与多模态 Agent 工作流")
    
    ai_col1 = [
        "**共情陪伴大模型 (LLM)**：",
        "- 选用 ModelScope 托管的 MiniMax-M2.5 共情陪伴大模型，生成温暖体贴的对话内容。",
        "- **Prompt Engineering**：系统 System Prompt 注入‘你是专业的心理健康助手心晴，擅长共情和安抚对话，回复带适当表情符号并控制在 200 字以内，给予鼓励。’",
        "- **打字机流式输出**：前端设计 typewriterEffect 算法，流式渲染，为低能量状态患者提供自然、缓和的倾听陪伴。"
    ]
    ai_col2 = [
        "**智能语音识别 (STT)**：",
        "- 选用 SiliconFlow 平台接入的 FunAudioLLM/SenseVoiceSmall 语音大模型。",
        "- **功能特性**：支持 App (mp3) 与 H5 (Blob WebM) 录音批量上传，具备低延迟、超高精确度识别性能，免除患者键盘打字的体力消耗，鼓励其以语音进行心理宣泄。"
    ]
    ai_col3 = [
        "**多模态 Agent 融合工作流**：",
        "- **HTP手绘视觉分析**：前端以 Stroke 算法量化笔触粗细、色彩冷暖、饱满度和偏心度，自动识别焦虑及抑郁风险趋势，向医生提供初筛报告。",
        "- **音视频情感分析**：上传音频触发 Python `main.py` 进行声学情感分析，视频触发 `emotion_http.py` 进行面部情緒推断，获取客观数据。",
        "- **近红外脑成像 (fNIRS) 映射**：患者线下 fNIRS 基准检测与线上手环指标、大模型情感度、主观量表交叉校验，实现全面评估。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "1. 共情 LLM 伴聊助手", ai_col1)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "2. 语音识别 (SenseVoiceSmall)", ai_col2)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "3. 多模态 Agent 评估机制", ai_col3)

    # ----------------------------------------------------
    # SLIDE 12: 四、商业模式及创新 - 商业模式（成本/收益/市场）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "四、商业模式", "商业模式、成本收益与市场前景分析")
    
    bm_c = [
        "**1. 研发与软硬件运维**：Spring Boot 云服务器托管、MySQL 数据库读写及 SSL 证书年费。",
        "**2. AI 服务调用费**：SiliconFlow (FunAudioLLM) 语音转文字与 ModelScope 大模型 API 调用流量成本。",
        "**3. 临床测试成本**：佳明手环采购/租赁成本、fNIRS 检测材料及受试者补偿费用。"
    ]
    bm_r = [
        "**1. B端机构系统集成与订阅**：面向高校心理测评中心、大中型企业工会心理援助板块提供 SaaS 系统年费订阅或专有云服务器定制化私有部署。",
        "**2. C端增值诊断与内容**：面向个人用户提供深度多模态身心报告导出服务，VIP 减压游戏扩展与冥想音频包订阅，以及在线对接专业医生进行视频问诊/挂号的抽成佣金。"
    ]
    bm_m = [
        "**1. 蓝海市场定位**：切入数字疗法 (Digital Therapeutics, DTx) 心理健康分支，填补纯硬件运动追踪同纯主观问卷 App 的空缺。",
        "**2. 核心市场壁垒**：打通‘运动生理-量表主观-大语言模型疗愈-临床医生监控’完整链条，依从性极强，具备很高的商业壁垒和品牌粘性。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "一、成本结构 (Costs)", bm_c)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "二、收益来源 (Revenue)", bm_r)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "三、市场定位 (Market)", bm_m)

    # ----------------------------------------------------
    # SLIDE 13: 四、商业模式及创新 - 行业落地价值（效率/测试用户数据）
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "四、商业模式", "行业落地价值与临床测试数据")
    
    val_l = [
        "**1. 临床入组测试闭环验证**：",
        "- 在**西北工业大学校医院**开展了为期 **15 天** 的抑郁及亚健康测试闭环。",
        "- **试验步骤**：线下入组筛选 -> 分配试验病例号 -> 签署知情同意书 -> 佳明 Connect 与手环一对一配对 -> 填写 SDS/PSS/SAS 初始基线量表 -> **近红外脑功能成像检测 (fNIRS)** -> 15 天连续佩戴监测与 EMA 日常状态填报 -> 第 15 天随访复测并收回手环。",
        "**2. 诊疗效率的大幅拉升**：",
        "- **手环生理客观指标**静默自动拉取同步，打消患者抗拒测评防御心理，提升患者依从率。",
        "- Web 后台量表及 HTP 房树人报告分类高亮警示（红/黄色），使得医生筛查重度风险患者的耗时缩减 **80% 以上**。"
    ]
    val_r = [
        "**1. 真实世界多模态生理心理库**：",
        "- 连续记录并结构化了入组用户 15天 完整的静息心率、深浅睡眠时段分布、血氧饱和度变化和 HRV 全天压力水平，积累了丰富的临床真实世界脱敏生理数据库 (RWD)。",
        "- 关联主观量表评测数据与 EMA 情绪起伏波形，为高校科研团队和精神科临床研究提供数据钻取分析支撑。",
        "**2. 多维度科学性交叉验证**：",
        "- 通过 15天 前后主观心理量表的分数改善，映射佳明物理压力指数的变化，并结合线下近红外脑成像 (fNIRS) 前后基准对照，提供了扎实的临床医学支撑，验证了心晴AI数字疗法的实用价值。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2), "1. 15天临床测试闭环与效率提升", val_l)
    add_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2), "2. 真实世界数据库与交叉验证价值", val_r)

    # ----------------------------------------------------
    # SLIDE 14: 四、商业模式及创新 - 公益计划
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, "四、商业模式", "社会公益价值与关怀计划")
    
    w1 = [
        "**高校危机干预防火墙**：",
        "- 联合西北工业大学等合作高校的心理辅导大厅及医院，无偿对被评估具有高危心理波动、中重度抑郁隐患的学生群体开放心晴AI的共情伴聊、跑酷减压游戏和冥想音频。",
        "- 当检测到高危聊天危机词汇时，静默向校医院后台发出红色警报，筑牢校园生命线第一道防线。"
    ]
    w2 = [
        "**公益手环租借援助计划**：",
        "- 联合公益基金会和红十字会，针对低保家庭、失业困难群体中的中重度抑郁症患者，免费提供智能健康手环的租借使用服务。",
        "- 免费将其健康数据链接入心晴AI，进行无感生理日常监控及恢复趋势跟踪，打消因贫困放弃心理监护的窘境。"
    ]
    w3 = [
        "**社区心理筛查与宣讲**：",
        "- 依托 h.playe.top H5 页面无需安装的极速传播优势，在社区服务中心、文化中心开展科普宣讲。",
        "- 开展免费的房树人手绘心理测验与 SCL-90 大众精神体检筛查，提升家庭对青少年情绪波动及老人睡眠健康状况的科学关爱意识。"
    ]
    
    add_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.2), "校园心理干预防线", w1)
    add_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.2), "公益手环租赁援助", w2)
    add_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.2), "社区心理科普筛查", w3)

    # Save to file
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "show.pptx"))
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
